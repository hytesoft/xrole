import json
import requests
import hashlib
import numpy as np
from learning.url_fingerprint import FingerprintDB
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
import logging
import time
import os
import subprocess
import sys
try:
    from tqdm import tqdm
except ImportError:
    tqdm = lambda x, **kwargs: x  # 没装tqdm时降级为普通迭代

def cosine_sim(a, b):
    a = np.array(a, dtype=np.float32)
    b = np.array(b, dtype=np.float32)
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b)))

# 日志配置
root_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
log_dir = os.path.join(root_dir, 'logs')
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, 'fetch_and_learn.log')
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s %(levelname)s %(message)s',
    handlers=[
        logging.FileHandler(log_file, encoding='utf-8'),
        logging.StreamHandler()
    ]
)

def notify_exception(msg):
    # 这里可扩展为邮件/钉钉/企业微信等通知
    logging.error(f"[异常通知] {msg}")
    # 可扩展：os.system('curl ...')

def run_audio2text():
    """自动调用音视频转写脚本"""
    from configparser import ConfigParser
    conf_path = os.path.join(os.path.dirname(__file__), '../config/xrole.conf')
    import json
    with open(conf_path, 'r', encoding='utf-8') as f:
        conf = json.load(f)
    material_dir = conf.get('material_dir', '/data/xrole_materials')
    script_path = os.path.join(os.path.dirname(__file__), 'audio2text.py')
    subprocess.run([sys.executable, script_path, material_dir])

def fetch_and_learn(config_path="config/xrole.conf", weblist_path="config/weblists.json"):
    # 读取配置
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            config = json.load(f)
    except Exception as e:
        notify_exception(f"读取配置失败: {e}")
        return
    # 先获取当前 prompt 和 weblists
    xrole_role=config.get("role", "")
    prompt = config.get("prompt", "")
    try:
        with open(weblist_path, "r", encoding="utf-8") as f:
            weblists_data = json.load(f)
            weblists = weblists_data["weblists"]
    except Exception as e:
        notify_exception(f"读取 weblists.json 失败: {e}")
        return
    # 组织大模型输入
    current_urls = [item["url"] for item in weblists]
    llm_input = f"{prompt}\n\n这是一个{xrole_role}，他需要持续的学习新的专业知识，请检索你大脑中与他的需求内容高度相关的内容，统计出相关系数高的前几名网址,没有合适的就不要勉强。不要通用性的网址，不要首页、导航页、聚合页、门户页，也不要需要再点开二级页面才能看到内容的网址。只要专栏、专题、技术博客、学习列表等，打开网址就能直接看到系统化的学习内容或文章列表，方便后续数据采集，严禁编造或拼凑网址。每条请附简要说明。当前已抓取资源有：{current_urls}。请以JSON数组格式返回（每项为url和简要说明），如：[{{'url': '...', 'desc': '...'}}, ...]。只返回JSON数组，不要多余解释。"
    #llm_input = f"{prompt}\n\n请只推荐真实存在且权威的专业网站，优先推荐国际公认的学术数据库、行业协会、标准组织、权威学术期刊、政府/高校/研究机构官网，避免仅推荐商业公司或通用门户网站（如 LinkedIn、YouTube、O'Reilly、BMC 等）。推荐前请根据你大脑中的内容来源有多少来自该网址作为推荐依据，严禁编造或拼凑网址。每条请附简要说明。当前已抓取资源有：{current_urls}。请以JSON数组格式返回建议的下载资源列表（每项为url和简要说明），如：[{{'url': '...', 'desc': '...'}}, ...]。只返回JSON数组，不要多余解释。"
    llm_conf = config.get("llm", {})
    try:
        resp = requests.post(
            llm_conf.get("base_url"),
            headers={"Authorization": f"Bearer {llm_conf.get('api_key', '')}"},
            json={
                "model": llm_conf.get("model"),
                "messages": [{"role": "user", "content": llm_input}]
            },
            timeout=60
        )
        llm_text = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        print(llm_text)
        # 健壮提取大模型返回的JSON数组，防止思维过程干扰
        import re
        import ast
        match = re.search(r'(\[.*?\])', llm_text, re.DOTALL)
        if match:
            try:
                new_weblists = ast.literal_eval(match.group(1))
                # 兼容格式，补全 desc 字段
                for item in new_weblists:
                    if "desc" not in item:
                        item["desc"] = ""
                # 更新 weblists.json
                with open(weblist_path, "w", encoding="utf-8") as f:
                    json.dump({"weblists": new_weblists}, f, ensure_ascii=False, indent=2)
                logging.info(f"已根据大模型建议更新 weblists.json，资源数: {len(new_weblists)}")
                weblists = new_weblists
            except Exception as e:
                logging.error(f"大模型输出无法解析为JSON数组，原始输出: {llm_text}")
                notify_exception(f"大模型输出无法解析为JSON数组: {e}")
        else:
            logging.warning("大模型未返回有效的资源列表，继续用原有 weblists")
    except Exception as e:
        notify_exception(f"大模型生成资源列表失败: {e}")
    # 支持多 embedding 模型
    embedding_models = config.get("embedding_models")
    if not embedding_models:
        # 兼容老配置
        embedding_models = [{"name": "paraphrase-multilingual-MiniLM-L12-v2"}]
    embedder_dict = {}
    for m in embedding_models:
        model_name = m["name"]
        # 如果不是绝对路径，拼成绝对路径（修正：去掉多余的 ..）
        if not os.path.isabs(model_name):
            local_model_path = os.path.join(root_dir, "models", "sentence-transformers", model_name)
            local_model_path = os.path.abspath(local_model_path)
        else:
            local_model_path = model_name
        try:
            logging.info(f"尝试加载 embedding 模型: {local_model_path}")
            embedder_dict[model_name] = SentenceTransformer(local_model_path, local_files_only=True)
            logging.info(f"embedding 模型 {model_name} 已用本地路径 {local_model_path} 加载成功")
        except Exception as e:
            notify_exception(f"加载 embedding 模型 {model_name} 失败: {e}")
    if __name__ == "__main__":
        import argparse
        parser = argparse.ArgumentParser(description="xrole fetch_and_learn")
        parser.add_argument('--mode', choices=['all', 'network', 'local'], default='all', help='all=全部，network=只采集网络，local=只导入本地资料')
        args = parser.parse_args()

        # collection 支持
        collections = config.get("collections")
        if not collections:
            collections = ["xrole_docs"]
        # 初始化依赖
        fingerprint_db = FingerprintDB()
        qdrant_conf = config.get("qdrant", {})
        qdrant_client = QdrantClient(
            url=qdrant_conf.get("url"),
            api_key=qdrant_conf.get("api_key")
        )
        # 读取 weblists.json 和大模型推荐后，再获取 spider_url
        spider_conf = config.get("sprider", {})
        spider_url = spider_conf.get("url")
        if not spider_url:
            logging.warning("未配置 spider.url，跳过抓取")
            return
        if args.mode in ('all', 'network'):
            # 网络采集主流程
            # 步骤1：遍历 weblists，抓取A，指纹比对
            for item in weblists:
                url = item["url"].rstrip("/")
                if url is None:
                    continue
                try:
                    resp = requests.post(spider_url.rstrip("/") + "/fetch", json={"url": url}, timeout=30, verify=False)
                    resp.raise_for_status()
                    content_a = resp.json().get("content", "")
                    print(f"{url}")
                    print(content_a)
                    if not content_a:
                        logging.warning(f"{url} 抓取无内容（A阶段）")
                        continue
                    model_name = item.get("embedding_model") or embedding_models[0]["name"]
                    embedder = embedder_dict.get(model_name)
                    if not embedder:
                        logging.error(f"未找到 embedding 模型 {model_name}，跳过 {url}")
                        continue
                    vector_a = embedder.encode(content_a)
                    meta = {"url": url, "embedding_model": model_name}
                    # 统一查重+入库
                    insert_if_not_exists(url, content_a, vector_a, meta, fingerprint_db, qdrant_client, item.get("collection") or collections[0])
                    # 步骤2：A为新内容，发给大模型让其抽取url列表
                    llm_input2 = f"{prompt}\n\n请从以下内容中提取所有有价值的资源url，按时间顺序输出JSON数组（如: ['url1', 'url2', ...]），只返回数组，不要解释。内容：{content_a[:1000]}..."
                    try:
                        resp2 = requests.post(
                            llm_conf.get("base_url"),
                            headers={"Authorization": f"Bearer {llm_conf.get('api_key', '')}"},
                            json={
                                "model": llm_conf.get("model"),
                                "messages": [{"role": "user", "content": llm_input2}]
                            },
                            timeout=60
                        )
                        llm_text2 = resp2.json().get("choices", [{}])[0].get("message", {}).get("content", "")
                        import re
                        import ast
                        match2 = re.search(r'\[.*\]', llm_text2, re.DOTALL)
                        if match2:
                            url_list = ast.literal_eval(match2.group(0))
                            logging.info(f"大模型抽取到{len(url_list)}个url，准备抓取B阶段")
                        else:
                            logging.warning(f"大模型未返回有效url列表，跳过 {url}")
                            continue
                        # 步骤3：遍历大模型返回的url列表，逐个抓取B，指纹比对，遇到老数据提前终止
                        for b_url in url_list:
                            try:
                                resp_b = requests.post(spider_url.rstrip("/") + "/fetch", json={"url": b_url}, timeout=30, verify=False)
                                resp_b.raise_for_status()
                                content_b = resp_b.json().get("content", "")
                                print(content_b)
                                if not content_b:
                                    logging.warning(f"{b_url} 抓取无内容（B阶段）")
                                    continue
                                vector_b = embedder.encode(content_b)
                                exists_b = False
                                for _, fp in fingerprint_db.get_all_fingerprints():
                                    if cosine_sim(vector_b, fp) > 0.95:
                                        exists_b = True
                                        break
                                if exists_b:
                                    logging.info(f"{b_url} 内容B已存在，后续url跳过")
                                    break  # 按时间顺序，遇到老数据提前终止
                                # 新内容，入库
                                fingerprint_db.add_fingerprint(b_url, vector_b)
                                qdrant_client.upsert(
                                    collection_name=item.get("collection") or collections[0],
                                    points=[{
                                        "id": hashlib.md5(b_url.encode()).hexdigest(),
                                        "vector": vector_b.tolist(),
                                        "payload": {"url": b_url, "content": content_b, "embedding_model": model_name}
                                    }]
                                )
                                logging.info(f"{b_url} 新内容已入库 [模型:{model_name} 集合:{item.get('collection') or collections[0]}]")
                            except Exception as e:
                                logging.error(f"{b_url} B阶段抓取或入库失败: {e}")
                    except Exception as e:
                        notify_exception(f"大模型抽取url失败: {e}")
                except Exception as e:
                    logging.error(f"{url} A阶段抓取或大模型处理失败: {e}")
        if args.mode in ('all', 'local'):
            # 本地资料导入
            material_dir = config.get("material_dir", "/data/xrole_materials")
            import_materials(material_dir, embedder_dict, embedding_models, collections, fingerprint_db, qdrant_client)
    # 本地资料导入独立调用
    material_dir = config.get("material_dir", "/data/xrole_materials")
    import_materials(material_dir, embedder_dict, embedding_models, collections, fingerprint_db, qdrant_client)

def file_fingerprint(file_path, content):
    return hashlib.md5((file_path + content).encode('utf-8')).hexdigest()

def insert_if_not_exists(unique_id, content, vector, meta, fingerprint_db, qdrant_client, collection, sim_threshold=0.95):
    """
    统一内容去重+入库接口。
    unique_id: 可为url、file_path等唯一标识
    content: 原始文本内容
    vector: embedding向量
    meta: dict，附加元信息
    fingerprint_db: 指纹数据库实例
    qdrant_client: Qdrant实例
    collection: 向量库集合名
    sim_threshold: embedding相似度去重阈值
    """
    # 1. 内容hash查重
    content_hash = hashlib.md5(content.encode('utf-8')).hexdigest()
    for _, old_fp in fingerprint_db.get_all_fingerprints():
        if content_hash == old_fp:
            logging.info(f"内容hash已存在，跳过: {unique_id}")
            return False
    # 2. embedding相似度查重
    for _, old_fp in fingerprint_db.get_all_fingerprints():
        try:
            if isinstance(old_fp, (list, np.ndarray)) and cosine_sim(vector, old_fp) > sim_threshold:
                logging.info(f"内容embedding相似，跳过: {unique_id}")
                return False
        except Exception:
            continue
    # 3. 入库
    fingerprint_db.add_fingerprint(unique_id, content_hash)
    qdrant_client.upsert(
        collection_name=collection,
        points=[{
            "id": hashlib.md5(unique_id.encode()).hexdigest(),
            "vector": vector.tolist(),
            "payload": {**meta, "content_hash": content_hash}
        }]
    )
    logging.info(f"新内容已入库: {unique_id}")
    return True

def import_materials(material_dir, embedder_dict, embedding_models, collections, fingerprint_db, qdrant_client):
    """自动导入宿主机挂载的学习资料（去重，避免重复学习）"""
    supported_exts = [".txt", ".md", ".pdf", ".ppt", ".pptx"]
    try:
        import glob
        from pathlib import Path
        # 可选依赖
        try:
            import pdfplumber
        except ImportError:
            pdfplumber = None
        try:
            from pptx import Presentation
        except ImportError:
            Presentation = None
        files = []
        for ext in supported_exts:
            files.extend(glob.glob(os.path.join(material_dir, f"**/*{ext}"), recursive=True))
        for file_path in tqdm(files, desc="本地资料导入进度"):
            ext = Path(file_path).suffix.lower()
            content = ""
            if ext in [".txt", ".md"]:
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        content = f.read()
                except Exception as e:
                    logging.warning(f"读取 {file_path} 失败: {e}")
            elif ext == ".pdf" and pdfplumber:
                try:
                    with pdfplumber.open(file_path) as pdf:
                        content = "\n".join(page.extract_text() or '' for page in pdf.pages)
                except Exception as e:
                    logging.warning(f"解析 PDF {file_path} 失败: {e}")
            elif ext in [".ppt", ".pptx"] and Presentation:
                try:
                    prs = Presentation(file_path)
                    slides = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slides.append(shape.text)
                    content = "\n".join(slides)
                except Exception as e:
                    logging.warning(f"解析 PPT {file_path} 失败: {e}")
            if content:
                model_name = embedding_models[0]["name"]
                embedder = embedder_dict.get(model_name)
                if embedder:
                    vector = embedder.encode(content)
                    url = f"file://{file_path}"
                    meta = {"url": url, "embedding_model": model_name, "source": "material_dir"}
                    insert_if_not_exists(url, content, vector, meta, fingerprint_db, qdrant_client, collections[0])
                    logging.info(f"已导入学习资料: {file_path}")
    except Exception as e:
        logging.error(f"自动导入学习资料失败: {e}")
